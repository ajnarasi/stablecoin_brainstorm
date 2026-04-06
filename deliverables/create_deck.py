#!/usr/bin/env python3
"""
Fiserv Digital Pay: Crypto & Stablecoin Acceptance Strategy
Investor Day Presentation - May 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Fiserv brand colors
NAVY = RGBColor(0x00, 0x3B, 0x5C)       # #003B5C - primary dark
ORANGE = RGBColor(0xFF, 0x66, 0x00)      # #FF6600 - accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MED_GRAY = RGBColor(0x9B, 0x9B, 0x9B)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
TEAL = RGBColor(0x00, 0x80, 0x90)        # accent 2
GREEN = RGBColor(0x2E, 0x7D, 0x32)       # success
RED = RGBColor(0xC6, 0x28, 0x28)         # risk

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=14, default_color=DARK_GRAY):
    """lines: list of (text, size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = align
        p.space_after = Pt(4)
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None, header_color=NAVY, header_text_color=WHITE):
    """data: list of lists. First row is header."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.color.rgb = header_text_color
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = DARK_GRAY

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY)

# Orange accent bar at top
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

# Title
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Fiserv Digital Pay", 48, WHITE, True, PP_ALIGN.LEFT, "Calibri")
add_text_box(slide, Inches(1), Inches(2.9), Inches(11), Inches(0.8),
             "Crypto & Stablecoin Acceptance Strategy", 28, ORANGE, False, PP_ALIGN.LEFT, "Calibri")

# Subtitle components
add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
             "CommerceHub  |  Finxact  |  FIUSD  |  INDX", 18, WHITE, False, PP_ALIGN.LEFT)

# Bottom bar
add_shape(slide, Inches(0), Inches(6.8), W, Inches(0.7), RGBColor(0x00, 0x2A, 0x42))
add_text_box(slide, Inches(1), Inches(6.9), Inches(5), Inches(0.4),
             "May 2026  |  CONFIDENTIAL", 12, MED_GRAY, False)
add_text_box(slide, Inches(7), Inches(6.9), Inches(5.5), Inches(0.4),
             "Investor Day Strategy Presentation", 12, MED_GRAY, False, PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 2: EXECUTIVE SUMMARY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Executive Summary", 36, NAVY, True)

# Four stat cards
stats = [
    ("$2-4T", "Stablecoin market\nby 2030", NAVY),
    ("<1%", "Transaction cost\nvs 2-3% interchange", ORANGE),
    ("6M", "Merchant locations\nready for activation", TEAL),
    ("ONLY", "Vertically integrated\ncrypto payment stack", GREEN),
]

for i, (big, desc, color) in enumerate(stats):
    x = Inches(0.8 + i * 3.1)
    card = add_rounded_rect(slide, x, Inches(1.5), Inches(2.8), Inches(1.8), LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.3), Inches(1.65), Inches(2.2), Inches(0.8),
                 big, 42, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(2.4), Inches(2.2), Inches(0.7),
                 desc, 13, DARK_GRAY, False, PP_ALIGN.CENTER)

# Key messages
msgs = [
    "GENIUS Act (July 2025) legalized stablecoins as payment instruments in the US",
    "75% of NRF 2026 retailers implementing or planning agentic commerce",
    "Fiserv is a founding member of the x402 Foundation (April 2, 2026)",
    "Cross-border payments represent a $190 trillion annual market opportunity",
    "FIUSD is live in production -- bank-embedded stablecoin on Solana via Paxos/Circle",
]

for i, msg in enumerate(msgs):
    y = Inches(3.7 + i * 0.42)
    add_shape(slide, Inches(1.0), y + Inches(0.08), Inches(0.12), Inches(0.12), ORANGE)
    add_text_box(slide, Inches(1.3), y, Inches(11), Inches(0.4), msg, 13, DARK_GRAY)


# ============================================================
# SLIDE 3: MARKET OPPORTUNITY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Market Opportunity", 36, NAVY, True)

# Left side: Market sizing
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(6), Inches(0.5),
             "Stablecoin & Agentic Commerce Markets", 20, NAVY, True)

market_data = [
    ["Segment", "Current", "2030 Projected", "CAGR"],
    ["Stablecoin Market Cap", "$260B", "$2-4T", "~50%"],
    ["Cross-Border Payments", "$190T/yr", "$250T/yr", "~6%"],
    ["Agentic Commerce", "~$0", "$1-3T", "N/A (new)"],
    ["US Stablecoin Share", "3% of USD payments", "10% by 2031", "~25%"],
]
add_table(slide, Inches(0.8), Inches(2.0), Inches(6), Inches(2.5), 5, 4, market_data)

# Right side: Key drivers
add_text_box(slide, Inches(7.5), Inches(1.4), Inches(5), Inches(0.5),
             "Market Drivers", 20, NAVY, True)

drivers = [
    ("Regulatory Clarity", "GENIUS Act + MiCA create legal framework for stablecoins as payment instruments"),
    ("Cost Advantage", "Sub-1% stablecoin costs vs 2-3% card interchange eliminates ~$50B in merchant fees"),
    ("Speed", "Real-time settlement vs T+1 to T+3 on card rails -- cash flow unlocked"),
    ("AI Agent Demand", "McKinsey: US B2C agentic commerce could reach $1T by 2030"),
]

for i, (title, desc) in enumerate(drivers):
    y = Inches(2.1 + i * 1.1)
    add_rounded_rect(slide, Inches(7.5), y, Inches(5.2), Inches(0.95), LIGHT_GRAY)
    add_text_box(slide, Inches(7.8), y + Inches(0.08), Inches(4.6), Inches(0.35),
                 title, 13, ORANGE, True)
    add_text_box(slide, Inches(7.8), y + Inches(0.4), Inches(4.6), Inches(0.5),
                 desc, 11, DARK_GRAY)

# Bottom callout
add_rounded_rect(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2), NAVY)
add_text_box(slide, Inches(1.2), Inches(5.95), Inches(11), Inches(0.9),
             "\"Stablecoin payments will capture 10% of US dollar payments by 2031.\nThe question is not IF -- it's WHO captures the orchestration layer.\"", 16, WHITE, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4: COMPETITIVE LANDSCAPE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Competitive Landscape", 36, NAVY, True)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
             "Only Fiserv has all four pillars of the vertically integrated stablecoin payment stack", 14, MED_GRAY)

comp_data = [
    ["Competitor", "Banking Core", "Own Stablecoin", "Settlement Network", "Merchant Acceptance", "Full Stack?"],
    ["Fiserv", "Finxact", "FIUSD", "INDX (1,100+ banks)", "CommerceHub + Clover (6M)", "YES"],
    ["Stripe", "No", "No (Bridge/USDC)", "Via Bridge ($1.1B acq)", "Stripe Checkout", "No"],
    ["Adyen", "No", "No", "No", "Adyen platform", "No"],
    ["Worldpay", "No", "No", "Via BVNK", "Worldpay platform", "No"],
    ["PayPal", "No", "PYUSD", "No", "PayPal/Venmo", "No"],
    ["Block/Square", "No", "No (BTC only)", "No", "Square POS", "No"],
]

col_w = [Inches(1.5), Inches(1.8), Inches(2.0), Inches(2.5), Inches(3.0), Inches(1.5)]
add_table(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(3.8), 7, 6, comp_data, col_w)

# Bottom insight
add_rounded_rect(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2), NAVY)
add_multi_text(slide, Inches(1.2), Inches(5.95), Inches(11), Inches(0.9), [
    ("Key Insight: Fiserv's moat is structural, not feature-based", 16, ORANGE, True, PP_ALIGN.CENTER),
    ("No competitor can replicate the closed loop: bank issues FIUSD > consumer pays merchant > INDX settles in real-time > merchant receives USD instantly", 13, WHITE, False, PP_ALIGN.CENTER),
])


# ============================================================
# SLIDE 5: VERTICALLY INTEGRATED STACK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Our Vertically Integrated Stack", 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "The only company with bank-issued stablecoin accepted at 6M merchants with real-time FDIC-insured settlement", 14, ORANGE)

# Stack visualization - 4 connected boxes
stack_items = [
    ("Finxact", "Cloud-Native Banking Core", "Multi-asset ledger\nAPI-first architecture\n10,000+ FI clients", Inches(0.8)),
    ("FIUSD", "Bank-Embedded Stablecoin", "Live on Solana\nPaxos/Circle backed\n1:1 USD peg", Inches(3.9)),
    ("INDX", "Real-Time Settlement", "1,100+ FDIC banks\n$25M FDIC insurance\n24/7/365 settlement", Inches(7.0)),
    ("CommerceHub / Clover", "Merchant Acceptance", "6M merchant locations\nE-comm + in-store\nSMB to enterprise", Inches(10.1)),
]

for (title, subtitle, desc, x) in stack_items:
    # Card background
    card = add_rounded_rect(slide, x, Inches(1.8), Inches(2.8), Inches(3.8), RGBColor(0x00, 0x4D, 0x73))
    # Title
    add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(2.4), Inches(0.5),
                 title, 20, ORANGE, True, PP_ALIGN.CENTER)
    # Subtitle
    add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(2.4), Inches(0.4),
                 subtitle, 12, WHITE, True, PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, x + Inches(0.2), Inches(3.1), Inches(2.4), Inches(2.0),
                 desc, 12, RGBColor(0xCC, 0xCC, 0xCC), False, PP_ALIGN.CENTER)

# Arrows between cards
for i in range(3):
    x = Inches(3.6 + i * 3.1)
    add_text_box(slide, x, Inches(3.2), Inches(0.5), Inches(0.5),
                 "\u2192", 28, ORANGE, True, PP_ALIGN.CENTER)

# Bottom flow
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8),
             "Bank issues FIUSD  \u2192  Consumer pays merchant  \u2192  INDX settles in real-time  \u2192  Merchant receives USD instantly",
             15, WHITE, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6: FLYWHEEL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "The FIUSD Flywheel", 36, NAVY, True)

# Flywheel steps as connected cards in a circle-ish layout
steps = [
    ("1", "FIs Issue FIUSD", "10,000+ bank clients\nissue via Finxact", Inches(4.5), Inches(1.5)),
    ("2", "Merchants Accept", "Cheaper than cards\n(sub-1% vs 2-3%)", Inches(8.5), Inches(2.5)),
    ("3", "Consumers Adopt", "Hold FIUSD in\nbank wallets", Inches(8.5), Inches(4.5)),
    ("4", "Volume Grows", "More FIs want in\ncosts decrease", Inches(4.5), Inches(4.5)),
    ("5", "Cycle Accelerates", "Network effects\ncompound", Inches(0.8), Inches(3.5)),
]

for (num, title, desc, x, y) in steps:
    card = add_rounded_rect(slide, x, y, Inches(3.2), Inches(1.3), LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(0.4), Inches(0.4),
                 num, 16, WHITE, True, PP_ALIGN.CENTER)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.1), y + Inches(0.1), Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()
    add_text_box(slide, x + Inches(0.1), y + Inches(0.1), Inches(0.4), Inches(0.4),
                 num, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.55), y + Inches(0.1), Inches(2.5), Inches(0.35),
                 title, 14, NAVY, True)
    add_text_box(slide, x + Inches(0.55), y + Inches(0.45), Inches(2.5), Inches(0.7),
                 desc, 11, DARK_GRAY)

# Seed strategy callout
add_rounded_rect(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(1.0), NAVY)
add_multi_text(slide, Inches(1.2), Inches(6.1), Inches(11), Inches(0.8), [
    ("Seed Strategy: Use Fiserv's own merchant settlement disbursements as initial FIUSD volume", 14, ORANGE, True, PP_ALIGN.LEFT),
    ("No dependency on consumer adoption. Fiserv controls enough of the system to prime the pump itself.", 12, WHITE, False, PP_ALIGN.LEFT),
])


# ============================================================
# SLIDE 7: WALLET PROVIDER ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Provider Architecture", 36, NAVY, True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "Layered approach behind Fiserv-owned abstraction layer -- merchants see \"Fiserv Digital Pay\", never the underlying providers", 13, MED_GRAY)

# Top layer - Fiserv Digital Pay
add_rounded_rect(slide, Inches(3.5), Inches(1.6), Inches(6.3), Inches(0.8), ORANGE)
add_text_box(slide, Inches(3.5), Inches(1.7), Inches(6.3), Inches(0.6),
             "\"Fiserv Digital Pay\"  --  Merchant-Facing Brand", 16, WHITE, True, PP_ALIGN.CENTER)

# Middle layer - CommerceHub orchestration
add_rounded_rect(slide, Inches(2.5), Inches(2.7), Inches(8.3), Inches(0.7), TEAL)
add_text_box(slide, Inches(2.5), Inches(2.8), Inches(8.3), Inches(0.5),
             "CommerceHub Crypto Orchestration Layer", 14, WHITE, True, PP_ALIGN.CENTER)

# Provider cards
providers = [
    ("Coinbase WaaS + CDP", "PRIMARY", "White-label MPC wallets\nx402 protocol creator\n$370B custody, SOC 2\n700+ wallets via WalletConnect", Inches(1.0)),
    ("Mesh Pay", "REACH", "300+ wallet/exchange integrations\nSmartFunding (patent-pending)\nApple Pay bridge\n45+ countries via Shift4", Inches(5.2)),
    ("Finxact + FIUSD + INDX", "OWNED FALLBACK", "Native wallet ledger\nFIUSD stablecoin\n$25M FDIC insurance\nClover in-store path", Inches(9.4)),
]

for (name, role, desc, x) in providers:
    add_rounded_rect(slide, x, Inches(3.7), Inches(3.5), Inches(2.8), LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.2), Inches(3.8), Inches(3.1), Inches(0.4),
                 name, 14, NAVY, True, PP_ALIGN.CENTER)
    # Role badge
    badge = add_rounded_rect(slide, x + Inches(0.8), Inches(4.2), Inches(1.9), Inches(0.35), ORANGE)
    add_text_box(slide, x + Inches(0.8), Inches(4.22), Inches(1.9), Inches(0.3),
                 role, 10, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(4.65), Inches(3.1), Inches(1.6),
                 desc, 11, DARK_GRAY, False, PP_ALIGN.CENTER)

# Design principles
add_text_box(slide, Inches(0.8), Inches(6.7), Inches(12), Inches(0.4),
             "Design Principles:  No provider >40% volume  |  90-day termination clauses  |  Card-rail fallback on every transaction  |  Circuit breaker on depeg",
             11, MED_GRAY, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 8: PROVIDER COMPARISON
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Provider Scoring: 7 Requirements", 36, NAVY, True)

score_data = [
    ["Requirement", "Coinbase", "Mesh", "Flexa*", "Heron*"],
    ["1. Crypto pay-in (global)", "9/10", "9/10", "7/10", "1/10"],
    ["2. Stablecoin pay-in (global)", "9/10", "8/10", "7/10", "3/10"],
    ["3. E-commerce + in-store", "5/10", "5/10", "7/10", "1/10"],
    ["4. Stablecoin disbursement", "8/10", "4/10", "2/10", "5/10"],
    ["5. Wallet storage on Finxact", "9/10", "4/10", "2/10", "2/10"],
    ["6. Stablecoin-to-fiat", "9/10", "8/10", "8/10", "4/10"],
    ["7. Any crypto-to-fiat (global)", "7/10", "9/10", "8/10", "1/10"],
    ["TOTAL", "56/70", "47/70", "41/70", "17/70"],
]

col_w2 = [Inches(3.2), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
add_table(slide, Inches(0.6), Inches(1.3), Inches(11.2), Inches(4.5), 9, 5, score_data, col_w2)

# Removal notes
add_rounded_rect(slide, Inches(0.8), Inches(6.0), Inches(5.5), Inches(1.1), RGBColor(0xFF, 0xEB, 0xEE))
add_multi_text(slide, Inches(1.1), Inches(6.1), Inches(5), Inches(0.9), [
    ("*Flexa Removed:", 12, RED, True, PP_ALIGN.LEFT),
    ("AMP token lost 98% of value. Only $14.1M raised after 7+ years. US/Canada only. No Clover integration confirmed.", 10, DARK_GRAY, False, PP_ALIGN.LEFT),
])

add_rounded_rect(slide, Inches(6.8), Inches(6.0), Inches(5.7), Inches(1.1), RGBColor(0xFF, 0xEB, 0xEE))
add_multi_text(slide, Inches(7.1), Inches(6.1), Inches(5.2), Inches(0.9), [
    ("*Heron Removed:", 12, RED, True, PP_ALIGN.LEFT),
    ("Zero verifiable funding, team, customers, or certifications. Three different companies share the name. Cannot be evaluated.", 10, DARK_GRAY, False, PP_ALIGN.LEFT),
])


# ============================================================
# SLIDE 9: SECURITY COMPARISON
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Security & Compliance Comparison", 36, NAVY, True)

sec_data = [
    ["Dimension", "Coinbase", "Mesh", "Flexa", "Heron"],
    ["SOC 2 Type II", "Yes (Deloitte)", "Yes", "No", "No"],
    ["SOC 1 Type II", "Yes (Deloitte)", "No", "No", "No"],
    ["Insurance", "$320M Lloyd's", "None", "AMP collateral", "None"],
    ["Assets Custody", "$370B+", "$0 (non-custodial)", "N/A", "$0"],
    ["US Licenses", "49 states + OCC", "None confirmed", "FinCEN MSB + MTLs", "None"],
    ["Int'l Licenses", "EU MiCA, SG, JP", "None confirmed", "None", "None"],
    ["KYC/AML", "Full internal", "Delegated", "Full internal", "Unknown"],
    ["Security Incidents", "1 (no crypto lost)", "None known", "None known", "Unknown"],
]

add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.8), 9, 5, sec_data,
          [Inches(2.2), Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)])

add_rounded_rect(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8), LIGHT_GRAY)
add_text_box(slide, Inches(1.1), Inches(6.4), Inches(11.2), Inches(0.6),
             "Coinbase is the only provider with both SOC 1 and SOC 2 Type II certifications audited by a Big Four firm (Deloitte), plus $320M insurance via Lloyd's of London.",
             12, DARK_GRAY, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10: AGENTIC COMMERCE PROTOCOLS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Agentic Commerce Protocol Landscape", 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "Five complementary protocols forming the agentic payment stack -- Fiserv is positioned across all of them", 13, ORANGE)

protocols = [
    ("x402", "Linux Foundation\n(Coinbase-created)", "HTTP-native\nstablecoin settlement", "Fiserv is\nFOUNDING MEMBER"),
    ("MPP", "Stripe / Tempo", "Session-based\nbatched payments", "Support via\nCommerceHub"),
    ("AP2", "Google\n+ 60 partners", "Agent authorization\n& trust", "Ecosystem\nparticipant"),
    ("TAP", "Visa", "Cryptographic\nagent identity", "Fiserv partnership\nannounced"),
    ("Agent Pay", "Mastercard", "Card infrastructure\nfor agents", "Expanded Fiserv\npartnership"),
]

for i, (name, creator, purpose, fiserv_role) in enumerate(protocols):
    x = Inches(0.5 + i * 2.55)
    card = add_rounded_rect(slide, x, Inches(1.7), Inches(2.35), Inches(3.5), RGBColor(0x00, 0x4D, 0x73))
    add_text_box(slide, x + Inches(0.15), Inches(1.85), Inches(2.05), Inches(0.5),
                 name, 22, ORANGE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(2.4), Inches(2.05), Inches(0.6),
                 creator, 11, MED_GRAY, False, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(3.1), Inches(2.05), Inches(0.6),
                 purpose, 11, WHITE, False, PP_ALIGN.CENTER)
    # Fiserv role badge
    add_rounded_rect(slide, x + Inches(0.15), Inches(3.9), Inches(2.05), Inches(0.6), ORANGE)
    add_text_box(slide, x + Inches(0.15), Inches(3.95), Inches(2.05), Inches(0.5),
                 fiserv_role, 10, WHITE, True, PP_ALIGN.CENTER)

# Bottom: CommerceHub Agent Pay
add_rounded_rect(slide, Inches(1.5), Inches(5.6), Inches(10.3), Inches(1.3), TEAL)
add_multi_text(slide, Inches(2.0), Inches(5.7), Inches(9.3), Inches(1.1), [
    ("CommerceHub Agent Pay: Protocol-Agnostic Orchestration", 18, WHITE, True, PP_ALIGN.CENTER),
    ("One CommerceHub endpoint. Detects whether the agent speaks x402, MPP, AP2, or TAP. Routes accordingly.", 13, WHITE, False, PP_ALIGN.CENTER),
    ("Defaults to FIUSD settlement. Bank-verified identity via Finxact. USD guaranteed via INDX.", 12, RGBColor(0xCC, 0xCC, 0xCC), False, PP_ALIGN.CENTER),
])


# ============================================================
# SLIDE 11: MPP BUILD VS IMPLEMENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Build vs. Implement: Protocol Strategy", 36, NAVY, True)

# Left: Don't Build
add_rounded_rect(slide, Inches(0.8), Inches(1.4), Inches(5.8), Inches(5.2), LIGHT_GRAY)
add_text_box(slide, Inches(1.0), Inches(1.5), Inches(5.4), Inches(0.5),
             "Don't Build a Proprietary Protocol", 18, RED, True)

dont_build = [
    "Fiserv publicly committed to x402 Foundation (Apr 2, 2026)",
    "Adyen (closest analog) is NOT building their own",
    "Protocol adoption requires ecosystem (22+ members vs 0)",
    "12-18 months engineering cost for zero network effects",
    "Every Fiserv advantage deliverable as implementation layer",
]
for i, point in enumerate(dont_build):
    add_shape(slide, Inches(1.2), Inches(2.15 + i * 0.55), Inches(0.12), Inches(0.12), RED)
    add_text_box(slide, Inches(1.5), Inches(2.05 + i * 0.55), Inches(4.8), Inches(0.5),
                 point, 12, DARK_GRAY)

# Right: Do Build This
add_rounded_rect(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(5.2), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.4), Inches(0.5),
             "Build the Best Implementation Layer", 18, GREEN, True)

do_build = [
    "Protocol-agnostic gateway (x402 + MPP + AP2 + TAP)",
    "Default to FIUSD settlement on Solana",
    "Bank-verified agent identity via Finxact KYC",
    "6M merchant auto-enablement via software update",
    "Fiserv Agent SDK (TypeScript + Python)",
    "\"Fiserv Commerce Profile\" within x402 specs",
]
for i, point in enumerate(do_build):
    add_shape(slide, Inches(7.4), Inches(2.15 + i * 0.55), Inches(0.12), Inches(0.12), GREEN)
    add_text_box(slide, Inches(7.7), Inches(2.05 + i * 0.55), Inches(4.8), Inches(0.5),
                 point, 12, DARK_GRAY)

# Bottom: strategic pattern
add_rounded_rect(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5), NAVY)
add_text_box(slide, Inches(1.2), Inches(6.85), Inches(11), Inches(0.4),
             "Same strategic pattern that made CommerceHub successful: don't own the network, own the orchestration", 13, WHITE, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 12: x402 DISTRIBUTION ADVANTAGE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "x402 Distribution Advantage", 36, WHITE, True)

# Big stat
add_text_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.5),
             "One Software Update =\n6 Million Agent-Payable Merchants", 40, ORANGE, True, PP_ALIGN.CENTER)

# Comparison
add_rounded_rect(slide, Inches(1.0), Inches(3.5), Inches(5.2), Inches(2.0), RGBColor(0x00, 0x4D, 0x73))
add_text_box(slide, Inches(1.3), Inches(3.7), Inches(4.6), Inches(0.4),
             "Stripe's Approach", 18, RED, True)
add_text_box(slide, Inches(1.3), Inches(4.2), Inches(4.6), Inches(1.0),
             "Convince merchants one by one\nEach integration requires developer effort\n100+ services at MPP launch (March 2026)", 13, WHITE)

add_rounded_rect(slide, Inches(7.1), Inches(3.5), Inches(5.2), Inches(2.0), RGBColor(0x00, 0x4D, 0x73))
add_text_box(slide, Inches(7.4), Inches(3.7), Inches(4.6), Inches(0.4),
             "Fiserv's Approach", 18, GREEN, True)
add_text_box(slide, Inches(7.4), Inches(4.2), Inches(4.6), Inches(1.0),
             "Push software update to Clover + CommerceHub\nZero merchant effort required\n6,000,000 merchants overnight", 13, WHITE)

add_text_box(slide, Inches(5.5), Inches(4.0), Inches(1.5), Inches(0.5),
             "VS", 24, ORANGE, True, PP_ALIGN.CENTER)

# Bottom
add_text_box(slide, Inches(1), Inches(6.0), Inches(11.3), Inches(0.8),
             "Every AI shopping agent that speaks x402 can immediately pay at every Fiserv merchant.\nFiserv becomes the largest x402 merchant network overnight.", 14, WHITE, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 13: BLUE OCEAN STRATEGY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Blue Ocean: Four Actions Framework", 36, NAVY, True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "The blue ocean is not \"crypto payments for merchants\" -- it's a bank-issued stablecoin payment network", 13, MED_GRAY)

actions = [
    ("ELIMINATE", "Crypto volatility risk for merchants\n\nAuto-convert everything to USD via INDX in real-time. Merchants never touch crypto.", RED, Inches(0.8), Inches(1.6)),
    ("REDUCE", "Integration complexity\n\nOne API toggle in CommerceHub. One button on Clover. Merchant does nothing.", RGBColor(0xE6, 0x5C, 0x00), Inches(7.0), Inches(1.6)),
    ("RAISE", "Regulatory confidence\n\nFDIC-insured settlement. Bank-embedded stablecoin. Regulated banking core. GENIUS Act compliant.", TEAL, Inches(0.8), Inches(4.0)),
    ("CREATE", "Bank-merchant stablecoin loop\n\nBanks earn 60-80bps (more than interchange). Merchants pay 80-100bps (less than cards). Card networks eliminated from the equation.", GREEN, Inches(7.0), Inches(4.0)),
]

for (title, desc, color, x, y) in actions:
    card = add_rounded_rect(slide, x, y, Inches(5.5), Inches(2.1), LIGHT_GRAY)
    # Color accent bar on left
    add_shape(slide, x, y, Inches(0.08), Inches(2.1), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(5.0), Inches(0.4),
                 title, 18, color, True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.55), Inches(5.0), Inches(1.4),
                 desc, 12, DARK_GRAY)

# Economics callout
add_rounded_rect(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.7), NAVY)
add_text_box(slide, Inches(1.2), Inches(6.48), Inches(11), Inches(0.5),
             "Result: Bank earns MORE per transaction, merchant pays LESS -- because the card network's share is eliminated", 14, ORANGE, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDES 14-17: PROTOTYPES (Problem + Solution + How It Works + Impact)
# ============================================================

def build_prototype_slide(prs, proto_num, title, problem_text, solution_text, steps, impacts, tagline):
    """Build a standardized prototype slide with Problem/Solution/How/Impact layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

    # Badge + Title
    add_rounded_rect(slide, Inches(0.8), Inches(0.3), Inches(1.8), Inches(0.35), ORANGE)
    add_text_box(slide, Inches(0.8), Inches(0.32), Inches(1.8), Inches(0.3),
                 f"PROTOTYPE {proto_num}", 11, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.8), Inches(0.25), Inches(9), Inches(0.5),
                 title, 28, NAVY, True)

    # Problem (left, red)
    add_rounded_rect(slide, Inches(0.6), Inches(0.9), Inches(6.0), Inches(1.3), RGBColor(0xFF, 0xEB, 0xEE))
    add_shape(slide, Inches(0.6), Inches(0.9), Inches(0.06), Inches(1.3), RED)
    add_multi_text(slide, Inches(0.9), Inches(0.95), Inches(5.5), Inches(1.2), [
        ("THE PROBLEM", 10, RED, True, PP_ALIGN.LEFT),
        (problem_text, 10, DARK_GRAY, False, PP_ALIGN.LEFT),
    ])

    # Solution (left, green)
    add_rounded_rect(slide, Inches(0.6), Inches(2.35), Inches(6.0), Inches(1.3), RGBColor(0xE8, 0xF5, 0xE9))
    add_shape(slide, Inches(0.6), Inches(2.35), Inches(0.06), Inches(1.3), GREEN)
    add_multi_text(slide, Inches(0.9), Inches(2.4), Inches(5.5), Inches(1.2), [
        ("THE SOLUTION", 10, GREEN, True, PP_ALIGN.LEFT),
        (solution_text, 10, DARK_GRAY, False, PP_ALIGN.LEFT),
    ])

    # How It Works (right column)
    add_text_box(slide, Inches(7.0), Inches(0.9), Inches(5.5), Inches(0.4),
                 "HOW IT WORKS", 10, NAVY, True)
    for i, (num, desc) in enumerate(steps):
        y = Inches(1.35 + i * 0.52)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0), y, Inches(0.3), Inches(0.3))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ORANGE
        circle.line.fill.background()
        add_text_box(slide, Inches(7.0), y + Inches(0.02), Inches(0.3), Inches(0.25),
                     num, 10, WHITE, True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(7.45), y, Inches(5.0), Inches(0.45),
                     desc, 9, DARK_GRAY)

    # Impact bar (bottom, navy)
    add_rounded_rect(slide, Inches(0.6), Inches(4.2), Inches(11.9), Inches(1.6), NAVY)
    add_text_box(slide, Inches(0.9), Inches(4.25), Inches(3), Inches(0.3),
                 "POTENTIAL IMPACT", 10, ORANGE, True)
    for i, (big, desc) in enumerate(impacts):
        x = Inches(0.8 + i * 3.05)
        add_text_box(slide, x, Inches(4.55), Inches(2.7), Inches(0.5),
                     big, 24, ORANGE, True, PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(5.05), Inches(2.7), Inches(0.45),
                     desc, 9, WHITE, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.4),
                 tagline, 11, WHITE, False, PP_ALIGN.CENTER)

    return slide


# SLIDE 14: Yield Sweep
build_prototype_slide(prs, 1, "Merchant Yield Sweep",
    problem_text="SMB merchants on Clover hold $15K-$50K in idle settlement balances earning zero interest. 99% of small businesses consume zero treasury management services. Across 6M merchant locations, billions sit dormant -- generating no yield for merchants and no revenue for Fiserv.",
    solution_text="AI-powered treasury agent auto-sweeps idle merchant balances into yield-bearing FIUSD positions on Finxact. ML model predicts cash needs using 180 days of CommerceHub transaction history. Configurable risk tolerance. Instant unsweep via INDX when cash is needed. One toggle in Clover to enable.",
    steps=[
        ("1", "Settlement data feeds into AI Treasury Agent from Clover/CommerceHub"),
        ("2", "ML model predicts 3-day outflows with seasonal adjustment (87% confidence)"),
        ("3", "Decision Gate validates safeguards: hard floor, gradual ramp, depeg check"),
        ("4", "Approved excess sweeps to FIUSD yield position on Finxact (4.2% APY)"),
        ("5", "INDX provides instant USD liquidity when merchant needs cash back"),
    ],
    impacts=[
        ("$847/mo", "Per merchant\nearnings on idle cash"),
        ("$750M/yr", "Revenue at full\npenetration (6M merchants)"),
        ("4.2% APY", "Yield on FIUSD\nvia Finxact"),
        ("Zero effort", "One toggle in Clover\nfully automated"),
    ],
    tagline="Fiserv becomes the ONLY payment processor that helps merchants earn money, not just spend it."
)

# SLIDE 15: Agent Pay
build_prototype_slide(prs, 2, "Pay-by-Agent x402 Commerce",
    problem_text="AI shopping agents (Claude, GPT) can browse and compare, but cannot buy. Card rails require human-in-the-loop authentication (3DS, CVV entry). This breaks the autonomous agentic flow entirely. Plus: 2-3% interchange fees and T+1 to T+3 settlement delays for merchants.",
    solution_text="CommerceHub exposes x402-compatible Agent Checkout endpoint. AI agents discover payment requirements, sign FIUSD payments (EIP-3009), and complete purchases in a single HTTP cycle -- under 3 seconds, fully autonomous. Bank-verified identity via Finxact KYC. Configurable spending guardrails.",
    steps=[
        ("1", "AI agent sends HTTP GET to CommerceHub merchant catalog"),
        ("2", "Gateway returns HTTP 402 Payment Required + payment instructions"),
        ("3", "Agent signs FIUSD payment via EIP-3009 on Solana/Base"),
        ("4", "Verifier checks signature, Finxact KYC tier, spending limits"),
        ("5", "Settler executes on-chain; INDX settles merchant in USD instantly"),
        ("6", "Agent receives cryptographic receipt -- all in under 3 seconds"),
    ],
    impacts=[
        ("6M", "Merchants agent-payable\nwith one software update"),
        ("<3 sec", "End-to-end purchase\nfully autonomous"),
        ("0.1%", "Agent payment fee\nvs 2.9% card interchange"),
        ("$150M/yr", "Revenue by 2030\nat scale"),
    ],
    tagline="Stripe convinces merchants one at a time. Fiserv flips a switch for 6 million."
)

# SLIDE 16: Supplier Pay
build_prototype_slide(prs, 3, "Instant Supplier Pay",
    problem_text="Restaurants on Clover place 15-30 supply orders/week. Each costs 2-3% in card interchange ($1,000-$1,500/month). Distributors wait 2-5 days for settlement. Early-pay discounts (2%/net-10) go uncaptured. A restaurant doing $50K/month in supplies loses $2,500/month in total waste.",
    solution_text="AI procurement agent monitors ingredient usage via Clover sales data, predicts reorder needs using ML, auto-generates POs grouped by supplier, pays instantly in FIUSD via CommerceHub. Supplier receives USD via INDX in 2.7 seconds. Early-pay discounts captured automatically.",
    steps=[
        ("1", "Clover POS tracks every sale; BOM maps menu items to ingredients"),
        ("2", "ML model predicts ingredient depletion (e.g., chicken in 2 days)"),
        ("3", "AI agent auto-generates POs grouped by supplier with MOQ enforcement"),
        ("4", "Early-pay discount applied (2% if paid within 24h vs net-30)"),
        ("5", "FIUSD payment via Finxact; supplier gets USD via INDX in 2.7 seconds"),
    ],
    impacts=[
        ("$1,200/mo", "Per merchant savings\n(fees + discounts)"),
        ("2.7 sec", "Supplier settlement\n(vs 2-5 days)"),
        ("0%", "Card interchange\non supply orders"),
        ("$690M/yr", "Revenue potential\nat full penetration"),
    ],
    tagline="Both merchant and supplier on Fiserv rails. Card fees eliminated. Discounts captured automatically."
)

# SLIDE 17: Cross-Border
build_prototype_slide(prs, 4, "Cross-Border Instant Settlement",
    problem_text="Cross-border card processing costs 3-5% in fees plus 2-4% FX markup. Settlement takes T+3 days through 4 intermediaries. A $1,000 MXN-to-USD transaction costs the merchant $47-$60 in fees. $6.3 trillion in cross-border e-commerce annually.",
    solution_text="CommerceHub auto-detects cross-border transactions (currency mismatch), converts local currency to FIUSD at real-time FX rates with 30-second rate lock, and settles merchant in USD via INDX in 3 seconds. Stablecoin is transport only -- merchant always receives dollars.",
    steps=[
        ("1", "International buyer pays in local currency (MXN, EUR, GBP)"),
        ("2", "CommerceHub detects currency mismatch, identifies corridor"),
        ("3", "FX engine locks rate for 30 seconds, converts to FIUSD on Solana"),
        ("4", "OFAC/sanctions screening passes (GENIUS Act compliant)"),
        ("5", "INDX converts FIUSD to USD, settles to merchant bank in 3 seconds"),
    ],
    impacts=[
        ("$55 saved", "Per $1,000 cross-border\ntransaction"),
        ("91.7%", "Cost reduction vs\ntraditional card rails"),
        ("3 sec", "Settlement time\n(vs 3 business days)"),
        ("$350M/yr", "Revenue at 70bps on\n$50B cross-border volume"),
    ],
    tagline="$6.3T cross-border e-commerce market. 90% cost reduction. Settlement in seconds, not days."
)



# ============================================================
# SLIDE 18: RISK ASSESSMENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Risk Assessment & Mitigation", 36, NAVY, True)

risk_data = [
    ["Risk", "Severity", "Probability", "Mitigation"],
    ["Regulatory reversal", "Catastrophic", "Low-Med", "FIUSD as deposit wrapper -- regulatory portable"],
    ["Stripe protocol lock-in", "High", "Medium", "6M auto-enable + dual protocol (x402 + MPP)"],
    ["Bank client resistance", "High", "High", "FIUSD = deposit + 60-80bps (> interchange)"],
    ["Consumer adoption stall", "Medium", "High", "Lead with B2B/agent use cases first"],
    ["Stablecoin depeg event", "Catastrophic", "Low", "Circuit breaker at 0.5% + instant conversion"],
    ["Internal BU misalignment", "High", "High", "C-suite sponsor + cross-BU task force"],
    ["Coinbase competitive conflict", "Medium", "Medium", "Abstraction layer + <40% volume cap"],
]

add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.5), 8, 4, risk_data,
          [Inches(2.8), Inches(1.5), Inches(1.5), Inches(6.5)])

# Bottom insight
add_rounded_rect(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(1.0), LIGHT_GRAY)
add_text_box(slide, Inches(1.2), Inches(6.1), Inches(11), Inches(0.8),
             "Every risk has a designed mitigation. The greatest risk is inaction -- competitors are moving NOW. Stripe acquired Bridge for $1.1B. Worldpay partnered with BVNK. Adyen joined every major protocol foundation.",
             12, DARK_GRAY, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 19: IMPLEMENTATION ROADMAP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Implementation Roadmap", 36, NAVY, True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "7-week parallel build  |  4 engineering teams  |  10-13 dedicated engineers  |  All prototypes demo-ready by May", 13, MED_GRAY)

roadmap_data = [
    ["Week", "Yield Sweep\n(2-3 engineers)", "Pay-by-Agent\n(3-4 engineers)", "Supplier Pay\n(2-3 engineers)", "Cross-Border\n(2-3 engineers)"],
    ["1", "Cash flow model", "x402 endpoint", "Inventory ingestion", "Detection engine"],
    ["2", "Finxact sweep API", "Payment verification", "AI procurement agent", "FX conversion"],
    ["3", "INDX liquidity", "Agent SDK (TS+Py)", "FIUSD B2B payment", "INDX integration"],
    ["4", "Dashboard widget", "INDX settlement", "Dashboard + tracker", "Comparison view"],
    ["5", "Demo polish", "Demo agent + dash", "Demo polish", "Demo polish"],
    ["6-7", "Buffer + rehearsal", "Buffer + rehearsal", "Buffer + rehearsal", "Buffer + rehearsal"],
]

add_table(slide, Inches(0.3), Inches(1.6), Inches(12.7), Inches(3.8), 7, 5, roadmap_data,
          [Inches(1.0), Inches(2.9), Inches(2.9), Inches(2.9), Inches(2.9)])

# MVP scope
add_text_box(slide, Inches(0.8), Inches(5.6), Inches(5.5), Inches(0.4),
             "MVP Scope Cuts for Timeline", 16, NAVY, True)

mvp_items = [
    "Yield Sweep: Simple rule-based (not ML) prediction",
    "Pay-by-Agent: FIUSD on Solana only, x402 only",
    "Supplier Pay: Single restaurant + single distributor",
    "Cross-Border: Single corridor (MXN -> USD)",
]
for i, item in enumerate(mvp_items):
    add_text_box(slide, Inches(1.0), Inches(6.1 + i * 0.32), Inches(5.3), Inches(0.3),
                 item, 11, DARK_GRAY)

# Presentation flow
add_text_box(slide, Inches(7.0), Inches(5.6), Inches(5.5), Inches(0.4),
             "Investor Day Flow (10 min)", 16, NAVY, True)

flow_items = [
    "(1 min) The Opportunity -- $2-4T market",
    "(1 min) Our Unique Position -- stack diagram",
    "(2 min) Live Demo: Yield Sweep",
    "(2 min) Live Demo: Pay-by-Agent",
    "(2 min) Live Demo: Supplier Pay + Cross-Border",
    "(1 min) Architecture + Ask",
]
for i, item in enumerate(flow_items):
    add_text_box(slide, Inches(7.2), Inches(6.1 + i * 0.28), Inches(5.3), Inches(0.28),
                 item, 10, DARK_GRAY)


# ============================================================
# SLIDE 20: THE ASK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "The Ask", 44, WHITE, True, PP_ALIGN.CENTER)

asks = [
    ("1", "Executive Sponsor", "C-suite leader to own the cross-BU outcome across CommerceHub, Finxact, INDX, and Clover"),
    ("2", "Cross-BU Task Force", "Dedicated team with P&L authority and representation from each business unit. 18-month mandate"),
    ("3", "Commercial Discussions", "Authorization to proceed with Coinbase and Mesh provider evaluations and commercial terms"),
    ("4", "Engineering Allocation", "10-13 dedicated engineers for 4 parallel prototype tracks through Investor Day and beyond"),
]

for i, (num, title, desc) in enumerate(asks):
    y = Inches(1.6 + i * 1.3)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()
    add_text_box(slide, Inches(1.5), y + Inches(0.05), Inches(0.6), Inches(0.5),
                 num, 22, WHITE, True, PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2.5), y + Inches(0.02), Inches(9), Inches(0.4),
                 title, 22, ORANGE, True)
    add_text_box(slide, Inches(2.5), y + Inches(0.45), Inches(9), Inches(0.5),
                 desc, 14, WHITE)

# Bottom quote
add_shape(slide, Inches(0), Inches(6.3), W, Inches(1.2), RGBColor(0x00, 0x2A, 0x42))
add_text_box(slide, Inches(1.5), Inches(6.45), Inches(10.3), Inches(0.8),
             "\"We are not just adding a payment method. We are building the infrastructure\nfor the next era of commerce.\"", 16, ORANGE, False, PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
output_path = "/Users/ajnarasi/Documents/Work/Projects/stablecoin/deliverables/Fiserv_Crypto_Strategy_Investor_Day.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
